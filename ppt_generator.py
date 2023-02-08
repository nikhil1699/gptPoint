# This is a test file to generate PPT using python uses hard coded instructions to
# generate an ugly ppt

import collections
import collections.abc
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt
from pptx.enum.dml import MSO_THEME_COLOR

title = 'Automated Presentation Creating Process'

pylogo = '2151056.jpeg'
pptlogo = '2151148.jpeg'
prs = Presentation()

slide = prs.slides.add_slide(prs.slide_layouts[6])
prs.slide_width = Inches(16)
prs.slide_height = Inches(9)

shape = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, 0, Inches(9 / 1.5), Inches(16), Inches(9 / 8.5)
)
shape.shadow.inherit = False
fill = shape.fill
fill.solid()
fill.fore_color.rgb = RGBColor(255, 0, 0)
shape.text = title
line = shape.line
line.color.rgb = RGBColor(255, 0, 0)
logo1 = slide.shapes.add_picture(pylogo, Inches(13.8), Inches(6.0), height=Inches(1.0), width=Inches(1.0))
logo2 = slide.shapes.add_picture(pptlogo, Inches(14.5), Inches(5.8), height=Inches(1.5), width=Inches(1.5))


slide = prs.slides.add_slide(prs.slide_layouts[6])

shape = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, 0, Inches(0.5), Inches(16), Inches(0.3))
shape.shadow.inherit = False
fill = shape.fill
fill.solid()
fill.fore_color.rgb = RGBColor(255, 0, 0)
shape.text = "How to Add a Chart"
line = shape.line
line.color.rgb = RGBColor(255, 0, 0)
logo1 = slide.shapes.add_picture(pylogo, Inches(14.5), Inches(0.4), height=Inches(0.5), width=Inches(0.5))
logo2 = slide.shapes.add_picture(pptlogo, Inches(15.0), Inches(0.4), height=Inches(0.5), width=Inches(0.5))

from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.chart.data import ChartData
from pptx.util import Inches
import numpy as np
import datetime

N = 100

random_x = np.random.randn(N) + 10
random_y = np.random.randn(N) + 5
random_z = np.random.randn(N) + 20

dte = datetime.datetime.today()
dt_lst = [dte - datetime.timedelta(days=i) for i in range(N)]

chart_data = ChartData()
chart_data.categories = dt_lst
chart_data.add_series('Data 1', random_x)
chart_data.add_series('Data 2', random_y)
chart_data.add_series('Data 3', random_z)

x, y, cx, cy = Inches(1), Inches(2), Inches(14), Inches(6)
chart = slide.shapes.add_chart(
    XL_CHART_TYPE.LINE, x, y, cx, cy, chart_data
).chart
chart.has_legend = True
chart.legend.include_in_layout = False
chart.series[2].smooth = True

slide = prs.slides.add_slide(prs.slide_layouts[6])

shape = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, 0, Inches(0.5), Inches(16), Inches(0.3))
shape.shadow.inherit = False
fill = shape.fill
fill.solid()
fill.fore_color.rgb = RGBColor(255, 0, 0)
shape.text = "Add an image damnit !"
line = shape.line
line.color.rgb = RGBColor(255, 0, 0)
logo1 = slide.shapes.add_picture(pylogo, Inches(14.5), Inches(0.4), height=Inches(0.5), width=Inches(0.5))
logo2 = slide.shapes.add_picture(pptlogo, Inches(15.0), Inches(0.4), height=Inches(0.5), width=Inches(0.5))

import plotly.graph_objects as go
import pandas as pd

df = pd.read_csv('https://raw.githubusercontent.com/plotly/datasets/\
718417069ead87650b90472464c7565dc8c2cb1c/sunburst-coffee-flavors-complete.csv')

fig = go.Figure(go.Sunburst(
    ids=df.ids,
    labels=df.labels,
    parents=df.parents))
fig.update_layout(uniformtext=dict(minsize=10, mode='hide'))

fig.write_image("img.png")

imgpth = 'img.png'

left = top = Inches(1)
pic = slide.shapes.add_picture(imgpth, left, top)


slide = prs.slides.add_slide(prs.slide_layouts[6])

shape = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, 0, Inches(0.5), Inches(16), Inches(0.3))
shape.shadow.inherit = False
fill = shape.fill
fill.solid()
fill.fore_color.rgb = RGBColor(255, 0, 0)
shape.text = "Wassup Beaches !"
line = shape.line
line.color.rgb = RGBColor(255, 0, 0)
logo1 = slide.shapes.add_picture(pylogo, Inches(14.5), Inches(0.4), height=Inches(0.5), width=Inches(0.5))
logo2 = slide.shapes.add_picture(pptlogo, Inches(15.0), Inches(0.4), height=Inches(0.5), width=Inches(0.5))

left = Inches(1)
top = Inches(2)
width = Inches(12)
height = Inches(5)

text_box = slide.shapes.add_textbox(left, top, width, height)

tb = text_box.text_frame
tb.text = 'Akkad bakkad Bambe Bo, Assi Nabbe poore Sau, Sau me nikla dhaaga,dum daba ke bhaaga'

prg = tb.add_paragraph()
prg.text = " "

prg = tb.add_paragraph()
prg.text = "They will find the Ring, and kill the one who carries it."


slide = prs.slides.add_slide(prs.slide_layouts[6])

shape = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, 0, Inches(4.0), Inches(16), Inches(1.0))
shape.shadow.inherit = False
fill = shape.fill
fill.solid()
fill.fore_color.rgb = RGBColor(255, 0, 0)
shape.text = "Thank You"
line = shape.line
line.color.rgb = RGBColor(255, 0, 0)
logo1 = slide.shapes.add_picture(pylogo, Inches(14.5), Inches(4.0), height=Inches(1.0), width=Inches(1.0))
logo2 = slide.shapes.add_picture(pptlogo, Inches(15.0), Inches(4.0), height=Inches(1.0), width=Inches(1.0))

prs.save('autogenerated_presentation.pptx')
